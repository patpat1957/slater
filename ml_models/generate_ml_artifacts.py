#!/usr/bin/env python3
"""
Generate ML-ready artifacts for Lottery Ball State Prediction Models.

Creates:
  1. ml_training_data.csv          - Example ML-ready spreadsheet
  2. ml_training_schema.json       - JSON schema for training data
  3. ml_database_schema.sql        - Database schema (PostgreSQL)
  4. ml_lottery_pipeline_deck.pptx - PowerPoint deck on ML pipeline
  5. ml_feature_config.json        - Feature engineering config

Ball States tracked per number per draw:
  - Winner ON/OFF: was this number drawn (1/0)?
  - Prediction ON/OFF: did our model predict this number (1/0)?
  - Best Practice ON/OFF: does this number satisfy strategy filters (1/0)?
"""

import csv
import json
import random
import datetime
import os

# ─────────────────────────────────────────────────────────────────────────────
# CONFIGURATION: Feature definitions from LotteryOptimizerEngine.js
# ─────────────────────────────────────────────────────────────────────────────

BALL_STATE_FEATURES = {
    "winner_state": {
        "description": "Was this ball number drawn in the actual draw?",
        "type": "binary",
        "values": {"ON": 1, "OFF": 0},
        "source": "actual draw results from API extraction"
    },
    "prediction_state": {
        "description": "Did our model predict this ball number for this draw?",
        "type": "binary",
        "values": {"ON": 1, "OFF": 0},
        "source": "model output / optimizer pred5 lines"
    },
    "best_practice_state": {
        "description": "Does this ball satisfy best-practice strategy filters?",
        "type": "binary",
        "values": {"ON": 1, "OFF": 0},
        "source": "checkCombination() pass/fail with 0 fails"
    }
}

FEATURE_COLUMNS = [
    # ── Identity
    ("draw_id", "string", "Unique draw identifier (state_game_date_drawtime)"),
    ("draw_date", "date", "Draw date YYYY-MM-DD"),
    ("state_code", "string", "2-letter US state code"),
    ("game_type", "string", "Game type: pick3/pick4/pick5/powerball/megamil/etc"),
    ("draw_time", "string", "midday or evening"),
    ("ball_position", "int", "Position index (0-based) in the draw"),
    ("ball_number", "int", "The actual ball number value"),

    # ── Ball State (TARGET variables + prediction tracking)
    ("winner_state", "binary", "1=drawn in actual result, 0=not drawn (TARGET for ML)"),
    ("prediction_state", "binary", "1=predicted by model, 0=not predicted"),
    ("best_practice_state", "binary", "1=passes all strategy filters, 0=fails one or more"),

    # ── Frequency Features (from computeFreqFromEntries)
    ("overall_freq", "float", "Raw frequency count across all history"),
    ("overall_freq_norm", "float", "Normalized frequency (0-1 scale, divided by max)"),
    ("temporal_decay_freq", "float", "Frequency with temporal decay lambda=0.02"),
    ("positional_freq", "float", "Frequency at this specific ball position"),
    ("positional_freq_norm", "float", "Normalized positional frequency (0-1)"),

    # ── Skip/Recency Features (from computeDetailedSkips)
    ("skip_value", "int", "Draws since last appearance (0=appeared last draw)"),
    ("skip_normalized", "float", "Skip / max_possible_skip (0-1)"),
    ("is_overdue", "binary", "1 if skip > SA threshold for this game"),
    ("sa_level", "int", "Which SA level (1-5) this number belongs to"),

    # ── Hot/Cold Classification (from getHCPL)
    ("is_hot", "binary", "1 if freq >= 60% of max frequency"),
    ("is_cold", "binary", "1 if freq < 30% of max frequency and freq > 0"),
    ("is_warm", "binary", "1 if between hot and cold thresholds"),
    ("is_prime_digit", "binary", "1 if last digit is 2,3,5,7"),
    ("is_low_digit", "binary", "1 if digit/last-digit is 0-4"),

    # ── Pair Co-occurrence (from scoreCombo pair analysis)
    ("avg_pair_cooccurrence", "float", "Average co-occurrence score with other drawn numbers"),
    ("max_pair_cooccurrence", "float", "Max co-occurrence with any single other number"),

    # ── FA Level (from computeFALevels)
    ("fa_level", "int", "Frequency Analysis level (1-4, 1=top 30%)"),
    ("in_fa_top30", "binary", "1 if in top 30% frequency"),
    ("in_fa_top50", "binary", "1 if in top 50% frequency"),

    # ── Combo-Level Features (per draw, same for all balls in that draw)
    ("combo_sum", "int", "Sum of all drawn numbers"),
    ("combo_sum_in_range", "binary", "1 if sum within 85th percentile historical range"),
    ("combo_odd_count", "int", "Number of odd numbers in the draw"),
    ("combo_even_count", "int", "Number of even numbers in the draw"),
    ("combo_oe_pattern", "string", "Odd/Even pattern e.g. '3O/2E'"),
    ("combo_oe_is_top2", "binary", "1 if O/E pattern matches top-2 historical"),
    ("combo_consec_sets", "int", "Number of consecutive number sets (e.g., 5-6 is 1 set)"),
    ("combo_consec_ok", "binary", "1 if consecutive sets <= 2 (PDF strategy)"),

    # ── AC Value (PDF Strategy 1)
    ("combo_ac_value", "int", "Arithmetic Complexity value"),
    ("combo_ac_in_range", "binary", "1 if AC within ideal range for this game"),

    # ── Gap Analysis (PDF Strategy 4)
    ("combo_gap_conforms", "binary", "1 if gaps between numbers match 98% historical profile"),
    ("min_gap", "int", "Minimum gap between consecutive sorted numbers"),
    ("max_gap", "int", "Maximum gap between consecutive sorted numbers"),
    ("avg_gap", "float", "Average gap between consecutive sorted numbers"),

    # ── SK123 Rule (PDF Strategy 5)
    ("combo_has_sk123", "binary", "1 if combo includes number from last 3 draws"),
    ("sk123_match_count", "int", "How many numbers from last 3 draws are in this combo"),

    # ── HCPL Ratio (PDF Strategy 6)
    ("combo_hot_count", "int", "Number of hot numbers in combo"),
    ("combo_cold_count", "int", "Number of cold numbers in combo"),
    ("combo_prime_count", "int", "Number of prime-digit numbers in combo"),
    ("combo_low_count", "int", "Number of low-digit numbers in combo"),
    ("combo_hcpl_ok", "binary", "1 if HCPL ratio within ideal historical range"),

    # ── Start/End Coverage (PDF Strategy 7)
    ("combo_start_num", "int", "First number in sorted combo"),
    ("combo_end_num", "int", "Last number in sorted combo"),
    ("combo_start_ok", "binary", "1 if start number in top-80% historical starts"),
    ("combo_end_ok", "binary", "1 if end number in top-80% historical ends"),
    ("combo_spread", "int", "End - Start (range of combo)"),

    # ── Pattern/Group (PDF Strategy 8)
    ("combo_pattern", "string", "L/M/H zone pattern e.g. 'L-L-M-H-H'"),
    ("combo_in_top27_pattern", "binary", "1 if pattern is among top 27 historical"),
    ("combo_group", "string", "Decade grouping e.g. '1-1-1-1-1'"),
    ("combo_in_top60_group", "binary", "1 if group is among top 60 historical"),

    # ── Game Status (PDF Strategy 10)
    ("game_is_normal", "binary", "1 if game in Normal status, 0 if Unstable"),
    ("game_unstable_count", "int", "Number of recent draws outside expected ranges"),

    # ── Walk-Forward Backtest
    ("wf_hit_rate", "float", "Walk-forward backtest hit rate (0-1)"),
    ("wf_hits", "int", "Walk-forward hits out of test window"),
    ("wf_total", "int", "Walk-forward test window size"),

    # ── Star Rating (composite quality score)
    ("star_rating", "int", "1-5 star rating from computeEnhancedStarRating"),

    # ── Prediction Tracker Learning
    ("tracker_weight", "float", "Adaptive learning weight from PredictionTracker"),
    ("tracker_hit_history", "float", "Historical hit rate for this number from tracker"),

    # ── Temporal Features
    ("day_of_week", "int", "0=Mon, 6=Sun"),
    ("week_of_year", "int", "1-52"),
    ("month", "int", "1-12"),
    ("draws_since_jackpot", "int", "Draws since last jackpot winner (if available)"),
]

# ─────────────────────────────────────────────────────────────────────────────
# 1. GENERATE EXAMPLE ML-READY CSV
# ─────────────────────────────────────────────────────────────────────────────

def generate_example_csv(filepath, num_draws=200):
    """Generate realistic example ML training data CSV."""
    headers = [col[0] for col in FEATURE_COLUMNS]

    # Game configs matching LotteryOptimizerEngine.js
    games = [
        {"game": "pick3", "state": "CA", "n": 3, "pool": 10, "is_num": False},
        {"game": "pick4", "state": "CA", "n": 4, "pool": 10, "is_num": False},
        {"game": "powerball", "state": "CA", "n": 5, "pool": 69, "is_num": True},
        {"game": "megamil", "state": "CA", "n": 5, "pool": 70, "is_num": True},
        {"game": "fantasy5", "state": "CA", "n": 5, "pool": 39, "is_num": True},
    ]

    rows = []
    base_date = datetime.date(2025, 1, 1)

    for game_cfg in games:
        game = game_cfg["game"]
        state = game_cfg["state"]
        n = game_cfg["n"]
        pool = game_cfg["pool"]
        is_num = game_cfg["is_num"]

        for draw_idx in range(num_draws):
            draw_date = base_date + datetime.timedelta(days=draw_idx)
            draw_time = random.choice(["midday", "evening"])
            draw_id = f"{state}_{game}_{draw_date}_{draw_time}"

            # Generate "actual" winning numbers
            if is_num:
                winning = sorted(random.sample(range(1, pool + 1), n))
            else:
                winning = [random.randint(0, pool - 1) for _ in range(n)]

            # Generate "predicted" numbers (simulate model output)
            if is_num:
                predicted = sorted(random.sample(range(1, pool + 1), n))
            else:
                predicted = [random.randint(0, pool - 1) for _ in range(n)]

            combo_sum = sum(winning)
            odd_count = sum(1 for x in winning if x % 2 != 0)
            even_count = n - odd_count
            oe_pattern = f"{odd_count}O/{even_count}E"

            # AC value
            if is_num and n >= 4:
                diffs = set()
                for i in range(len(winning)):
                    for j in range(i + 1, len(winning)):
                        diffs.add(abs(winning[j] - winning[i]))
                ac_value = len(diffs) - (n - 1)
            else:
                ac_value = 0

            # Gaps
            if is_num:
                gaps = [winning[i+1] - winning[i] for i in range(len(winning)-1)]
                min_gap = min(gaps) if gaps else 0
                max_gap = max(gaps) if gaps else 0
                avg_gap = round(sum(gaps) / len(gaps), 2) if gaps else 0
            else:
                min_gap = max_gap = 0
                avg_gap = 0.0

            # Consecutive sets
            consec_sets = 0
            in_set = False
            for i in range(1, len(winning)):
                if winning[i] - winning[i-1] == 1:
                    if not in_set:
                        consec_sets += 1
                        in_set = True
                else:
                    in_set = False

            # Per-ball features
            for pos_idx, ball_num in enumerate(winning):
                winner_state = 1  # This ball WAS drawn
                prediction_state = 1 if ball_num in predicted else 0
                best_practice_state = 1 if random.random() > 0.3 else 0

                row = {
                    "draw_id": draw_id,
                    "draw_date": str(draw_date),
                    "state_code": state,
                    "game_type": game,
                    "draw_time": draw_time,
                    "ball_position": pos_idx,
                    "ball_number": ball_num,
                    "winner_state": winner_state,
                    "prediction_state": prediction_state,
                    "best_practice_state": best_practice_state,
                    "overall_freq": random.randint(10, 200),
                    "overall_freq_norm": round(random.uniform(0.1, 1.0), 4),
                    "temporal_decay_freq": round(random.uniform(5, 180), 2),
                    "positional_freq": random.randint(2, 80),
                    "positional_freq_norm": round(random.uniform(0.05, 1.0), 4),
                    "skip_value": random.randint(0, 25),
                    "skip_normalized": round(random.uniform(0, 1), 4),
                    "is_overdue": 1 if random.random() > 0.8 else 0,
                    "sa_level": random.randint(1, 5),
                    "is_hot": 1 if random.random() > 0.5 else 0,
                    "is_cold": 1 if random.random() > 0.7 else 0,
                    "is_warm": 0,
                    "is_prime_digit": 1 if (ball_num % 10) in [2, 3, 5, 7] else 0,
                    "is_low_digit": 1 if (ball_num % 10) < 5 else 0,
                    "avg_pair_cooccurrence": round(random.uniform(0, 1), 4),
                    "max_pair_cooccurrence": round(random.uniform(0, 1), 4),
                    "fa_level": random.randint(1, 4),
                    "in_fa_top30": 1 if random.random() > 0.7 else 0,
                    "in_fa_top50": 1 if random.random() > 0.5 else 0,
                    "combo_sum": combo_sum,
                    "combo_sum_in_range": 1 if random.random() > 0.15 else 0,
                    "combo_odd_count": odd_count,
                    "combo_even_count": even_count,
                    "combo_oe_pattern": oe_pattern,
                    "combo_oe_is_top2": 1 if random.random() > 0.4 else 0,
                    "combo_consec_sets": consec_sets,
                    "combo_consec_ok": 1 if consec_sets <= 2 else 0,
                    "combo_ac_value": ac_value,
                    "combo_ac_in_range": 1 if (4 <= ac_value <= 6 and n == 5) or (7 <= ac_value <= 10 and n == 6) or not is_num else 0,
                    "combo_gap_conforms": 1 if random.random() > 0.2 else 0,
                    "min_gap": min_gap,
                    "max_gap": max_gap,
                    "avg_gap": avg_gap,
                    "combo_has_sk123": 1 if random.random() > 0.1 else 0,
                    "sk123_match_count": random.randint(0, min(3, n)),
                    "combo_hot_count": random.randint(0, n),
                    "combo_cold_count": random.randint(0, 2),
                    "combo_prime_count": random.randint(0, n),
                    "combo_low_count": random.randint(0, n),
                    "combo_hcpl_ok": 1 if random.random() > 0.3 else 0,
                    "combo_start_num": winning[0] if is_num else winning[0],
                    "combo_end_num": winning[-1] if is_num else winning[-1],
                    "combo_start_ok": 1 if random.random() > 0.2 else 0,
                    "combo_end_ok": 1 if random.random() > 0.2 else 0,
                    "combo_spread": winning[-1] - winning[0] if is_num else 0,
                    "combo_pattern": "-".join(["L" if x <= pool//3 else "M" if x <= 2*pool//3 else "H" for x in winning]) if is_num else "N/A",
                    "combo_in_top27_pattern": 1 if random.random() > 0.3 else 0,
                    "combo_group": "-".join([str(sum(1 for x in winning if d*10 < x <= (d+1)*10)) for d in range(pool//10 + 1)]) if is_num else "N/A",
                    "combo_in_top60_group": 1 if random.random() > 0.2 else 0,
                    "game_is_normal": 1 if random.random() > 0.2 else 0,
                    "game_unstable_count": random.randint(0, 5),
                    "wf_hit_rate": round(random.uniform(0, 0.6), 4),
                    "wf_hits": random.randint(0, 15),
                    "wf_total": 30,
                    "star_rating": random.randint(1, 5),
                    "tracker_weight": round(random.uniform(0.5, 2.0), 4),
                    "tracker_hit_history": round(random.uniform(0, 0.4), 4),
                    "day_of_week": draw_date.weekday(),
                    "week_of_year": draw_date.isocalendar()[1],
                    "month": draw_date.month,
                    "draws_since_jackpot": random.randint(0, 30),
                }
                row["is_warm"] = 1 if (row["is_hot"] == 0 and row["is_cold"] == 0) else 0
                rows.append(row)

            # Also add non-winning balls (negative examples) - sample some
            if is_num:
                non_winning = [x for x in range(1, pool + 1) if x not in winning]
            else:
                all_balls = list(range(pool))
                non_winning = [x for x in all_balls if x not in winning]

            for ball_num in random.sample(non_winning, min(n, len(non_winning))):
                row = dict(rows[-1])  # copy combo-level features
                row["ball_position"] = -1  # not in a winning position
                row["ball_number"] = ball_num
                row["winner_state"] = 0  # NOT drawn
                row["prediction_state"] = 1 if ball_num in predicted else 0
                row["best_practice_state"] = 1 if random.random() > 0.5 else 0
                row["overall_freq"] = random.randint(5, 150)
                row["overall_freq_norm"] = round(random.uniform(0.05, 0.8), 4)
                row["temporal_decay_freq"] = round(random.uniform(3, 130), 2)
                row["positional_freq"] = random.randint(1, 50)
                row["positional_freq_norm"] = round(random.uniform(0.02, 0.7), 4)
                row["skip_value"] = random.randint(1, 40)
                row["skip_normalized"] = round(random.uniform(0.1, 1), 4)
                row["is_overdue"] = 1 if random.random() > 0.6 else 0
                row["is_hot"] = 1 if random.random() > 0.6 else 0
                row["is_cold"] = 1 if random.random() > 0.5 else 0
                row["is_warm"] = 1 if (row["is_hot"] == 0 and row["is_cold"] == 0) else 0
                row["is_prime_digit"] = 1 if (ball_num % 10) in [2, 3, 5, 7] else 0
                row["is_low_digit"] = 1 if (ball_num % 10) < 5 else 0
                rows.append(row)

    with open(filepath, 'w', newline='') as f:
        writer = csv.DictWriter(f, fieldnames=headers)
        writer.writeheader()
        writer.writerows(rows)

    print(f"  CSV: {len(rows)} rows -> {filepath}")
    return len(rows)


# ─────────────────────────────────────────────────────────────────────────────
# 2. GENERATE JSON SCHEMA
# ─────────────────────────────────────────────────────────────────────────────

def generate_json_schema(filepath):
    """Generate comprehensive JSON schema for ML training data."""
    schema = {
        "$schema": "https://json-schema.org/draft/2020-12/schema",
        "title": "Lottery ML Training Data Schema",
        "description": "Schema for machine learning training data based on lottery ball states (Winner ON/OFF, Prediction ON/OFF, Best Practice ON/OFF) and 20 optimizer strategy features from LotteryOptimizerEngine.js",
        "version": "1.0.0",
        "type": "object",
        "properties": {
            "metadata": {
                "type": "object",
                "properties": {
                    "created_at": {"type": "string", "format": "date-time"},
                    "source": {"type": "string", "const": "Lotto Extraction API v1.1.0 + LotteryOptimizerEngine.js"},
                    "feature_count": {"type": "integer"},
                    "ball_states": {
                        "type": "object",
                        "description": "The three core ball states tracked per number per draw",
                        "properties": {
                            "winner_state": BALL_STATE_FEATURES["winner_state"],
                            "prediction_state": BALL_STATE_FEATURES["prediction_state"],
                            "best_practice_state": BALL_STATE_FEATURES["best_practice_state"],
                        }
                    },
                    "supported_games": {
                        "type": "array",
                        "items": {"type": "string"},
                        "examples": ["pick3", "pick4", "pick5", "powerball", "megamil", "fantasy5", "superlotto", "cash5", "lotto"]
                    },
                    "ml_targets": {
                        "type": "object",
                        "description": "Target variables for different ML tasks",
                        "properties": {
                            "classification": {
                                "type": "object",
                                "properties": {
                                    "primary": {"const": "winner_state", "description": "Binary: will this number be drawn?"},
                                    "secondary": {"const": "best_practice_state", "description": "Binary: does this number pass strategy filters?"}
                                }
                            },
                            "regression": {
                                "type": "object",
                                "properties": {
                                    "score": {"const": "combo_score", "description": "Continuous: optimizer score from scoreCombo()"},
                                    "wf_hit_rate": {"const": "wf_hit_rate", "description": "Continuous: walk-forward hit rate"}
                                }
                            },
                            "ranking": {
                                "type": "object",
                                "properties": {
                                    "star_rating": {"const": "star_rating", "description": "Ordinal 1-5: combo quality rating"}
                                }
                            }
                        }
                    }
                }
            },
            "feature_definitions": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "name": {"type": "string"},
                        "dtype": {"type": "string", "enum": ["string", "int", "float", "binary", "date"]},
                        "description": {"type": "string"},
                        "category": {"type": "string", "enum": [
                            "identity", "ball_state", "frequency", "skip_recency",
                            "hot_cold", "pair_cooccurrence", "frequency_analysis",
                            "combo_level", "ac_value", "gap_analysis", "sk123_rule",
                            "hcpl_ratio", "start_end", "pattern_group", "game_status",
                            "backtest", "star_rating", "tracker", "temporal"
                        ]}
                    }
                }
            },
            "training_samples": {
                "type": "array",
                "items": {
                    "type": "object",
                    "description": "One row per ball-number per draw",
                    "properties": {col[0]: {"type": "number" if col[1] in ("int", "float", "binary") else "string"} for col in FEATURE_COLUMNS}
                }
            },
            "model_configs": {
                "type": "object",
                "description": "Recommended ML model configurations",
                "properties": {
                    "xgboost": {
                        "type": "object",
                        "properties": {
                            "task": {"const": "binary_classification"},
                            "target": {"const": "winner_state"},
                            "params": {
                                "type": "object",
                                "properties": {
                                    "objective": {"const": "binary:logistic"},
                                    "eval_metric": {"const": "auc"},
                                    "max_depth": {"type": "integer", "default": 6},
                                    "learning_rate": {"type": "number", "default": 0.05},
                                    "n_estimators": {"type": "integer", "default": 500},
                                    "subsample": {"type": "number", "default": 0.8},
                                    "colsample_bytree": {"type": "number", "default": 0.8},
                                    "min_child_weight": {"type": "integer", "default": 5},
                                    "scale_pos_weight": {"type": "string", "default": "auto (n_negative/n_positive)"}
                                }
                            }
                        }
                    },
                    "random_forest": {
                        "type": "object",
                        "properties": {
                            "task": {"const": "binary_classification"},
                            "target": {"const": "winner_state"},
                            "params": {
                                "type": "object",
                                "properties": {
                                    "n_estimators": {"type": "integer", "default": 300},
                                    "max_depth": {"type": "integer", "default": 12},
                                    "min_samples_split": {"type": "integer", "default": 10},
                                    "min_samples_leaf": {"type": "integer", "default": 5},
                                    "max_features": {"const": "sqrt"},
                                    "class_weight": {"const": "balanced"}
                                }
                            }
                        }
                    },
                    "lightgbm": {
                        "type": "object",
                        "properties": {
                            "task": {"const": "binary_classification"},
                            "target": {"const": "winner_state"},
                            "params": {
                                "type": "object",
                                "properties": {
                                    "objective": {"const": "binary"},
                                    "metric": {"const": "auc"},
                                    "num_leaves": {"type": "integer", "default": 31},
                                    "learning_rate": {"type": "number", "default": 0.05},
                                    "n_estimators": {"type": "integer", "default": 500},
                                    "feature_fraction": {"type": "number", "default": 0.8},
                                    "bagging_fraction": {"type": "number", "default": 0.8},
                                    "is_unbalance": {"type": "boolean", "default": True}
                                }
                            }
                        }
                    },
                    "neural_network": {
                        "type": "object",
                        "properties": {
                            "task": {"const": "binary_classification"},
                            "target": {"const": "winner_state"},
                            "architecture": {
                                "type": "object",
                                "properties": {
                                    "type": {"const": "feedforward"},
                                    "layers": {"type": "array", "default": [128, 64, 32, 1]},
                                    "activation": {"const": "relu (hidden), sigmoid (output)"},
                                    "dropout": {"type": "number", "default": 0.3},
                                    "optimizer": {"const": "adam"},
                                    "loss": {"const": "binary_crossentropy"},
                                    "batch_size": {"type": "integer", "default": 256},
                                    "epochs": {"type": "integer", "default": 100}
                                }
                            }
                        }
                    }
                }
            }
        }
    }

    # Add category to each feature
    feature_defs = []
    categories = {
        "draw_id": "identity", "draw_date": "identity", "state_code": "identity",
        "game_type": "identity", "draw_time": "identity", "ball_position": "identity",
        "ball_number": "identity",
        "winner_state": "ball_state", "prediction_state": "ball_state", "best_practice_state": "ball_state",
        "overall_freq": "frequency", "overall_freq_norm": "frequency",
        "temporal_decay_freq": "frequency", "positional_freq": "frequency", "positional_freq_norm": "frequency",
        "skip_value": "skip_recency", "skip_normalized": "skip_recency", "is_overdue": "skip_recency", "sa_level": "skip_recency",
        "is_hot": "hot_cold", "is_cold": "hot_cold", "is_warm": "hot_cold",
        "is_prime_digit": "hot_cold", "is_low_digit": "hot_cold",
        "avg_pair_cooccurrence": "pair_cooccurrence", "max_pair_cooccurrence": "pair_cooccurrence",
        "fa_level": "frequency_analysis", "in_fa_top30": "frequency_analysis", "in_fa_top50": "frequency_analysis",
        "star_rating": "star_rating",
        "tracker_weight": "tracker", "tracker_hit_history": "tracker",
        "day_of_week": "temporal", "week_of_year": "temporal", "month": "temporal", "draws_since_jackpot": "temporal",
        "wf_hit_rate": "backtest", "wf_hits": "backtest", "wf_total": "backtest",
    }
    for col_name, col_type, col_desc in FEATURE_COLUMNS:
        cat = categories.get(col_name, "combo_level")
        if "ac_" in col_name: cat = "ac_value"
        elif "gap" in col_name: cat = "gap_analysis"
        elif "sk123" in col_name: cat = "sk123_rule"
        elif "hcpl" in col_name or col_name in ("combo_hot_count","combo_cold_count","combo_prime_count","combo_low_count"): cat = "hcpl_ratio"
        elif "start" in col_name or "end" in col_name or "spread" in col_name: cat = "start_end"
        elif "pattern" in col_name or "group" in col_name: cat = "pattern_group"
        elif "normal" in col_name or "unstable" in col_name: cat = "game_status"
        feature_defs.append({"name": col_name, "dtype": col_type, "description": col_desc, "category": cat})

    schema["properties"]["feature_definitions"]["items"]["examples"] = feature_defs[:3]

    with open(filepath, 'w') as f:
        json.dump(schema, f, indent=2, default=str)

    # Also save feature definitions separately
    feat_path = filepath.replace("schema", "features")
    with open(feat_path, 'w') as f:
        json.dump({"features": feature_defs, "total_features": len(feature_defs), "ball_states": BALL_STATE_FEATURES}, f, indent=2)

    print(f"  JSON schema: {filepath}")
    print(f"  JSON features: {feat_path}")


# ─────────────────────────────────────────────────────────────────────────────
# 3. GENERATE DATABASE SCHEMA (PostgreSQL)
# ─────────────────────────────────────────────────────────────────────────────

def generate_sql_schema(filepath):
    """Generate PostgreSQL database schema for lottery ML pipeline."""
    sql = """-- ═══════════════════════════════════════════════════════════════════════════
--  Lottery ML Pipeline — Database Schema (PostgreSQL)
--  Stores: raw draws, ball states, ML features, predictions, model results
--  Compatible with: XGBoost, Random Forest, LightGBM, Neural Networks
-- ═══════════════════════════════════════════════════════════════════════════

-- ── 1. Raw Draw Results (from Lotto Extraction API) ──────────────────────
CREATE TABLE IF NOT EXISTS draws (
    draw_id         TEXT PRIMARY KEY,           -- e.g. "CA_powerball_2025-05-27_evening"
    draw_date       DATE NOT NULL,
    state_code      CHAR(2) NOT NULL,
    game_type       VARCHAR(20) NOT NULL,       -- pick3/pick4/powerball/megamil/etc
    draw_time       VARCHAR(10) NOT NULL,       -- midday/evening
    numbers         INTEGER[] NOT NULL,         -- winning numbers array
    bonus_number    INTEGER,                    -- Powerball/Mega Ball (nullable)
    combo_sum       INTEGER NOT NULL,
    created_at      TIMESTAMPTZ DEFAULT NOW(),
    UNIQUE(state_code, game_type, draw_date, draw_time)
);

CREATE INDEX idx_draws_date ON draws(draw_date DESC);
CREATE INDEX idx_draws_game ON draws(state_code, game_type);

-- ── 2. Ball States (core tracking: Winner/Prediction/BestPractice ON/OFF) ──
CREATE TABLE IF NOT EXISTS ball_states (
    id              BIGSERIAL PRIMARY KEY,
    draw_id         TEXT REFERENCES draws(draw_id),
    ball_number     INTEGER NOT NULL,
    ball_position   INTEGER,                    -- position in draw (-1 if not drawn)

    -- ═══ THE THREE CORE BALL STATES ═══
    winner_state    SMALLINT NOT NULL DEFAULT 0,    -- 1=drawn, 0=not drawn
    prediction_state SMALLINT NOT NULL DEFAULT 0,   -- 1=predicted, 0=not predicted
    best_practice_state SMALLINT NOT NULL DEFAULT 0,-- 1=passes filters, 0=fails

    -- Derived tracking
    hit_type        VARCHAR(20),                -- 'exact_hit','partial_hit','miss','false_positive'
    created_at      TIMESTAMPTZ DEFAULT NOW()
);

CREATE INDEX idx_ball_states_draw ON ball_states(draw_id);
CREATE INDEX idx_ball_states_number ON ball_states(ball_number);
CREATE INDEX idx_ball_states_winner ON ball_states(winner_state);

-- ── 3. ML Feature Vectors (one row per ball per draw) ────────────────────
CREATE TABLE IF NOT EXISTS ml_features (
    id              BIGSERIAL PRIMARY KEY,
    draw_id         TEXT REFERENCES draws(draw_id),
    ball_number     INTEGER NOT NULL,

    -- Ball States (duplicated for fast queries)
    winner_state    SMALLINT NOT NULL DEFAULT 0,
    prediction_state SMALLINT NOT NULL DEFAULT 0,
    best_practice_state SMALLINT NOT NULL DEFAULT 0,

    -- Frequency features
    overall_freq         REAL,
    overall_freq_norm    REAL,
    temporal_decay_freq  REAL,
    positional_freq      REAL,
    positional_freq_norm REAL,

    -- Skip/Recency
    skip_value      INTEGER,
    skip_normalized REAL,
    is_overdue      SMALLINT DEFAULT 0,
    sa_level        SMALLINT,

    -- Hot/Cold/Prime/Low
    is_hot          SMALLINT DEFAULT 0,
    is_cold         SMALLINT DEFAULT 0,
    is_warm         SMALLINT DEFAULT 0,
    is_prime_digit  SMALLINT DEFAULT 0,
    is_low_digit    SMALLINT DEFAULT 0,

    -- Pair co-occurrence
    avg_pair_cooccurrence REAL,
    max_pair_cooccurrence REAL,

    -- FA level
    fa_level        SMALLINT,
    in_fa_top30     SMALLINT DEFAULT 0,
    in_fa_top50     SMALLINT DEFAULT 0,

    -- Combo-level features
    combo_sum            INTEGER,
    combo_sum_in_range   SMALLINT DEFAULT 0,
    combo_odd_count      SMALLINT,
    combo_even_count     SMALLINT,
    combo_oe_pattern     VARCHAR(10),
    combo_oe_is_top2     SMALLINT DEFAULT 0,
    combo_consec_sets    SMALLINT,
    combo_consec_ok      SMALLINT DEFAULT 0,

    -- AC Value (Strategy 1)
    combo_ac_value       SMALLINT,
    combo_ac_in_range    SMALLINT DEFAULT 0,

    -- Gap Analysis (Strategy 4)
    combo_gap_conforms   SMALLINT DEFAULT 0,
    min_gap              SMALLINT,
    max_gap              SMALLINT,
    avg_gap              REAL,

    -- SK123 (Strategy 5)
    combo_has_sk123      SMALLINT DEFAULT 0,
    sk123_match_count    SMALLINT,

    -- HCPL (Strategy 6)
    combo_hot_count      SMALLINT,
    combo_cold_count     SMALLINT,
    combo_prime_count    SMALLINT,
    combo_low_count      SMALLINT,
    combo_hcpl_ok        SMALLINT DEFAULT 0,

    -- Start/End (Strategy 7)
    combo_start_num      INTEGER,
    combo_end_num        INTEGER,
    combo_start_ok       SMALLINT DEFAULT 0,
    combo_end_ok         SMALLINT DEFAULT 0,
    combo_spread         INTEGER,

    -- Pattern/Group (Strategy 8)
    combo_pattern        VARCHAR(20),
    combo_in_top27_pattern SMALLINT DEFAULT 0,
    combo_group          VARCHAR(30),
    combo_in_top60_group SMALLINT DEFAULT 0,

    -- Game status (Strategy 10)
    game_is_normal       SMALLINT DEFAULT 1,
    game_unstable_count  SMALLINT DEFAULT 0,

    -- Walk-forward backtest
    wf_hit_rate          REAL,
    wf_hits              INTEGER,
    wf_total             INTEGER,

    -- Star rating
    star_rating          SMALLINT,

    -- Tracker learning
    tracker_weight       REAL DEFAULT 1.0,
    tracker_hit_history  REAL DEFAULT 0.0,

    -- Temporal
    day_of_week          SMALLINT,
    week_of_year         SMALLINT,
    month                SMALLINT,
    draws_since_jackpot  INTEGER,

    created_at           TIMESTAMPTZ DEFAULT NOW()
);

CREATE INDEX idx_ml_features_draw ON ml_features(draw_id);
CREATE INDEX idx_ml_features_ball ON ml_features(ball_number);
CREATE INDEX idx_ml_features_winner ON ml_features(winner_state);

-- ── 4. ML Models (track trained models and their performance) ────────────
CREATE TABLE IF NOT EXISTS ml_models (
    model_id        TEXT PRIMARY KEY,           -- e.g. "xgb_powerball_v3_20250527"
    model_name      VARCHAR(100) NOT NULL,      -- "XGBoost Powerball Classifier"
    model_type      VARCHAR(30) NOT NULL,       -- xgboost/random_forest/lightgbm/neural_net
    game_type       VARCHAR(20) NOT NULL,
    state_code      CHAR(2),
    target_column   VARCHAR(50) NOT NULL,       -- "winner_state"
    feature_columns TEXT[] NOT NULL,            -- list of feature column names used
    hyperparameters JSONB NOT NULL,             -- model hyperparameters
    training_rows   INTEGER NOT NULL,
    training_date_range TEXT,                   -- "2024-01-01 to 2025-05-01"
    metrics         JSONB NOT NULL,             -- {"auc": 0.72, "accuracy": 0.68, "f1": 0.55, ...}
    model_blob_path TEXT,                       -- path to saved model file
    created_at      TIMESTAMPTZ DEFAULT NOW(),
    is_active       BOOLEAN DEFAULT TRUE
);

-- ── 5. ML Predictions (model outputs for upcoming draws) ─────────────────
CREATE TABLE IF NOT EXISTS ml_predictions (
    id              BIGSERIAL PRIMARY KEY,
    model_id        TEXT REFERENCES ml_models(model_id),
    target_draw_date DATE NOT NULL,
    state_code      CHAR(2) NOT NULL,
    game_type       VARCHAR(20) NOT NULL,
    draw_time       VARCHAR(10) NOT NULL,
    ball_number     INTEGER NOT NULL,
    predicted_probability REAL NOT NULL,         -- 0.0 to 1.0
    predicted_class SMALLINT NOT NULL,           -- 0 or 1
    confidence_tier VARCHAR(10),                 -- "high"/"medium"/"low"
    prediction_rank INTEGER,                     -- rank within this prediction set
    created_at      TIMESTAMPTZ DEFAULT NOW()
);

CREATE INDEX idx_predictions_date ON ml_predictions(target_draw_date);
CREATE INDEX idx_predictions_model ON ml_predictions(model_id);

-- ── 6. Prediction Outcomes (compare predictions vs actual) ───────────────
CREATE TABLE IF NOT EXISTS prediction_outcomes (
    id              BIGSERIAL PRIMARY KEY,
    prediction_id   BIGINT REFERENCES ml_predictions(id),
    draw_id         TEXT REFERENCES draws(draw_id),
    was_correct     BOOLEAN NOT NULL,
    actual_state    SMALLINT NOT NULL,           -- actual winner_state
    predicted_prob  REAL NOT NULL,
    created_at      TIMESTAMPTZ DEFAULT NOW()
);

-- ── 7. Feature Importance (from trained models) ──────────────────────────
CREATE TABLE IF NOT EXISTS feature_importance (
    id              BIGSERIAL PRIMARY KEY,
    model_id        TEXT REFERENCES ml_models(model_id),
    feature_name    VARCHAR(50) NOT NULL,
    importance_score REAL NOT NULL,
    importance_rank INTEGER NOT NULL,
    importance_type VARCHAR(20) DEFAULT 'gain'   -- gain/weight/cover/shap
);

-- ── VIEWS ─────────────────────────────────────────────────────────────────

-- Ball state summary per number (for quick lookups)
CREATE OR REPLACE VIEW v_ball_state_summary AS
SELECT
    ball_number,
    d.game_type,
    d.state_code,
    COUNT(*) as total_draws,
    SUM(bs.winner_state) as times_won,
    SUM(bs.prediction_state) as times_predicted,
    SUM(bs.best_practice_state) as times_best_practice,
    ROUND(AVG(bs.winner_state)::numeric, 4) as win_rate,
    ROUND(AVG(bs.prediction_state)::numeric, 4) as prediction_rate,
    ROUND(AVG(bs.best_practice_state)::numeric, 4) as best_practice_rate,
    -- Precision: when we predicted, how often was it a winner?
    CASE WHEN SUM(bs.prediction_state) > 0
         THEN ROUND(SUM(CASE WHEN bs.winner_state=1 AND bs.prediction_state=1 THEN 1 ELSE 0 END)::numeric
                     / SUM(bs.prediction_state)::numeric, 4)
         ELSE 0 END as prediction_precision,
    -- Recall: of actual winners, how many did we predict?
    CASE WHEN SUM(bs.winner_state) > 0
         THEN ROUND(SUM(CASE WHEN bs.winner_state=1 AND bs.prediction_state=1 THEN 1 ELSE 0 END)::numeric
                     / SUM(bs.winner_state)::numeric, 4)
         ELSE 0 END as prediction_recall
FROM ball_states bs
JOIN draws d ON bs.draw_id = d.draw_id
GROUP BY ball_number, d.game_type, d.state_code;

-- Model performance comparison
CREATE OR REPLACE VIEW v_model_performance AS
SELECT
    m.model_id,
    m.model_name,
    m.model_type,
    m.game_type,
    (m.metrics->>'auc')::real as auc,
    (m.metrics->>'accuracy')::real as accuracy,
    (m.metrics->>'f1')::real as f1_score,
    (m.metrics->>'precision')::real as precision,
    (m.metrics->>'recall')::real as recall,
    m.training_rows,
    m.created_at,
    m.is_active
FROM ml_models m
ORDER BY (m.metrics->>'auc')::real DESC;

-- Training data export view (join everything for ML pipeline)
CREATE OR REPLACE VIEW v_ml_training_data AS
SELECT
    f.*,
    d.numbers as actual_numbers,
    d.bonus_number as actual_bonus
FROM ml_features f
JOIN draws d ON f.draw_id = d.draw_id
ORDER BY d.draw_date DESC, f.ball_number;
"""

    with open(filepath, 'w') as f:
        f.write(sql)

    print(f"  SQL schema: {filepath}")


# ─────────────────────────────────────────────────────────────────────────────
# 4. GENERATE POWERPOINT DECK
# ─────────────────────────────────────────────────────────────────────────────

def generate_pptx(filepath):
    """Generate PowerPoint deck explaining the ML lottery pipeline."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_slide(title_text, content_items=None, layout_idx=1):
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        for p in title.text_frame.paragraphs:
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
        if content_items and len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for item in content_items:
                p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(18)
                p.space_after = Pt(8)
                if item.startswith("  "):
                    p.level = 1
                    p.font.size = Pt(16)
        return slide

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Lottery Ball State ML Pipeline"
    for p in slide.shapes.title.text_frame.paragraphs:
        p.font.size = Pt(40)
        p.font.bold = True
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "XGBoost | Random Forest | LightGBM | Neural Networks\nBased on 20 Optimizer Strategies + Ball State Tracking"

    # Slide 2: The Three Ball States
    add_slide("The Three Core Ball States", [
        "Every ball number in every draw has THREE binary states:",
        "",
        "1. Winner State (ON/OFF) - Was this number actually drawn?",
        "  TARGET variable for ML models",
        "  Source: actual draw results from Lotto Extraction API",
        "",
        "2. Prediction State (ON/OFF) - Did our model predict this number?",
        "  Tracks model output for accuracy measurement",
        "  Source: model predictions / optimizer pred5 lines",
        "",
        "3. Best Practice State (ON/OFF) - Does it pass strategy filters?",
        "  Based on 20 optimizer strategies (AC, Skip, Gap, Pattern, etc.)",
        "  Source: checkCombination() with 0 filter failures",
    ])

    # Slide 3: Ball State Matrix
    add_slide("Ball State Matrix - Decision Framework", [
        "WIN=ON + PREDICT=ON + BP=ON  --> Perfect Hit (model + strategy aligned)",
        "WIN=ON + PREDICT=ON + BP=OFF --> Lucky Hit (model right, strategy missed)",
        "WIN=ON + PREDICT=OFF + BP=ON --> Missed Winner (strategy right, model missed)",
        "WIN=ON + PREDICT=OFF + BP=OFF -> Total Miss (neither caught it)",
        "",
        "WIN=OFF + PREDICT=ON + BP=ON --> False Positive (both agreed, wrong)",
        "WIN=OFF + PREDICT=ON + BP=OFF -> Model False Positive",
        "WIN=OFF + PREDICT=OFF + BP=ON --> Strategy False Positive",
        "WIN=OFF + PREDICT=OFF + BP=OFF -> True Negative (correctly excluded)",
        "",
        "Goal: Maximize Perfect Hits, minimize Total Misses",
        "Track precision/recall for both prediction and best-practice states",
    ])

    # Slide 4: Feature Engineering
    add_slide("Feature Engineering - 60+ Features from 20 Strategies", [
        "FREQUENCY: overall_freq, temporal_decay_freq, positional_freq",
        "SKIP/RECENCY: skip_value, is_overdue, sa_level (1-5)",
        "HOT/COLD: is_hot, is_cold, is_warm, is_prime_digit, is_low_digit",
        "PAIR CO-OCCURRENCE: avg/max pair frequency with other numbers",
        "AC VALUE: Arithmetic Complexity (Strategy 1, ideal 4-6 for 5-num games)",
        "GAP ANALYSIS: min/max/avg gap, 98% conformance (Strategy 4)",
        "SK123: includes number from last 3 draws (Strategy 5)",
        "HCPL RATIO: hot/cold/prime/low counts in combo (Strategy 6)",
        "START/END: first/last number coverage (Strategy 7)",
        "PATTERN/GROUP: L/M/H zone pattern, decade grouping (Strategy 8)",
        "GAME STATUS: normal vs unstable detection (Strategy 10)",
        "TEMPORAL: day_of_week, month, week_of_year, draws_since_jackpot",
    ])

    # Slide 5: XGBoost Configuration
    add_slide("Model 1: XGBoost (Recommended Primary)", [
        "Task: Binary Classification (winner_state = 0 or 1)",
        "",
        "Key Hyperparameters:",
        "  objective: binary:logistic",
        "  eval_metric: auc (area under ROC curve)",
        "  max_depth: 6 (prevent overfitting on random noise)",
        "  learning_rate: 0.05 (slow learning = better generalization)",
        "  n_estimators: 500 (with early stopping patience=50)",
        "  subsample: 0.8 (row sampling per tree)",
        "  colsample_bytree: 0.8 (feature sampling per tree)",
        "  scale_pos_weight: auto (handles class imbalance)",
        "",
        "Strengths: Handles missing values, feature importance ranking,",
        "  native support for imbalanced classes (lottery = very imbalanced)",
    ])

    # Slide 6: Random Forest
    add_slide("Model 2: Random Forest", [
        "Task: Binary Classification (winner_state = 0 or 1)",
        "",
        "Key Hyperparameters:",
        "  n_estimators: 300 (number of trees)",
        "  max_depth: 12 (deeper than XGBoost to capture complex interactions)",
        "  min_samples_split: 10",
        "  min_samples_leaf: 5",
        "  max_features: sqrt (random feature subset per split)",
        "  class_weight: balanced (auto-adjust for class imbalance)",
        "",
        "Strengths: Robust to overfitting, no hyperparameter sensitivity,",
        "  parallel training, good baseline model",
        "",
        "Use Case: Ball-level classification, feature importance comparison",
    ])

    # Slide 7: LightGBM
    add_slide("Model 3: LightGBM (Fastest Training)", [
        "Task: Binary Classification + Ranking",
        "",
        "Key Hyperparameters:",
        "  objective: binary",
        "  metric: auc",
        "  num_leaves: 31 (leaf-wise growth = faster than depth-wise)",
        "  learning_rate: 0.05",
        "  n_estimators: 500",
        "  feature_fraction: 0.8",
        "  bagging_fraction: 0.8",
        "  is_unbalance: true",
        "",
        "Strengths: 10x faster than XGBoost on large datasets,",
        "  native categorical feature support,",
        "  excellent for iterative prediction with frequent retraining",
    ])

    # Slide 8: Neural Network
    add_slide("Model 4: Neural Network (Deep Learning)", [
        "Architecture: Feedforward Network",
        "  Input Layer: 60+ features (normalized 0-1)",
        "  Hidden 1: 128 units + ReLU + Dropout(0.3)",
        "  Hidden 2: 64 units + ReLU + Dropout(0.3)",
        "  Hidden 3: 32 units + ReLU + Dropout(0.2)",
        "  Output: 1 unit + Sigmoid (probability 0-1)",
        "",
        "Training Config:",
        "  Optimizer: Adam (lr=0.001)",
        "  Loss: Binary Cross-Entropy",
        "  Batch Size: 256",
        "  Epochs: 100 with early stopping (patience=15)",
        "  Class weights: inversely proportional to frequency",
        "",
        "Alternative: LSTM for sequential draw pattern learning",
    ])

    # Slide 9: ML Pipeline Architecture
    add_slide("End-to-End ML Pipeline Architecture", [
        "Step 1: DATA EXTRACTION",
        "  Lotto Extraction API --> raw draw results (JSON/CSV)",
        "",
        "Step 2: FEATURE ENGINEERING",
        "  LotteryOptimizerEngine.js --> 60+ features per ball per draw",
        "  Ball states: winner_state, prediction_state, best_practice_state",
        "",
        "Step 3: TRAINING",
        "  PostgreSQL --> pandas DataFrame --> train/val/test split (70/15/15)",
        "  Train XGBoost + RF + LightGBM + NN (ensemble approach)",
        "",
        "Step 4: PREDICTION",
        "  For next draw: compute features --> ensemble predict --> rank balls",
        "",
        "Step 5: EVALUATION",
        "  Compare predictions vs actual --> update ball_states --> retrain",
    ])

    # Slide 10: Training/Validation Strategy
    add_slide("Training & Validation Strategy", [
        "CRITICAL: Time-series split (NO random shuffle!)",
        "  Train: draws from 2020-01-01 to 2024-12-31",
        "  Validation: draws from 2025-01-01 to 2025-03-31",
        "  Test: draws from 2025-04-01 to present",
        "",
        "Walk-Forward Validation:",
        "  Train on N draws --> predict draw N+1 --> evaluate",
        "  Slide window forward --> retrain --> repeat",
        "  Matches walkForwardBacktest() in LotteryOptimizerEngine.js",
        "",
        "Metrics to Track:",
        "  AUC-ROC: primary metric (handles imbalanced classes)",
        "  Precision@K: top-K predictions that are actual winners",
        "  Hit Rate: percentage of actual winners in our prediction set",
        "  Calibration: predicted probability vs actual win rate",
    ])

    # Slide 11: Class Imbalance
    add_slide("Handling Class Imbalance (Critical for Lottery)", [
        "Problem: In a 5/69 game (Powerball), each draw has:",
        "  5 winning balls (positive) vs 64 non-winning (negative)",
        "  Class ratio: ~7% positive / ~93% negative",
        "",
        "Solutions (use ALL of these):",
        "  1. scale_pos_weight (XGBoost) / class_weight (RF) = auto-balance",
        "  2. SMOTE oversampling of positive class during training",
        "  3. Threshold tuning: don't use 0.5 cutoff, optimize F1-score threshold",
        "  4. Ranking approach: rank all balls by probability, pick top-N",
        "  5. Focal Loss (Neural Net): down-weight easy negatives",
        "",
        "Evaluation: Use AUC-ROC and Precision@K, NOT accuracy",
        "  (a model that always predicts 'not drawn' gets 93% accuracy!)",
    ])

    # Slide 12: Feature Importance
    add_slide("Expected Feature Importance Ranking", [
        "Tier 1 - Strongest Predictors (from existing optimizer):",
        "  overall_freq_norm, temporal_decay_freq, skip_value",
        "  is_hot, sa_level, positional_freq_norm",
        "",
        "Tier 2 - Strategy Signals:",
        "  combo_ac_in_range, combo_gap_conforms, combo_has_sk123",
        "  combo_in_top27_pattern, combo_hcpl_ok",
        "  avg_pair_cooccurrence",
        "",
        "Tier 3 - Contextual:",
        "  game_is_normal, combo_sum_in_range, combo_oe_is_top2",
        "  tracker_weight, wf_hit_rate",
        "",
        "Tier 4 - Temporal (weak but informative):",
        "  day_of_week, month, draws_since_jackpot",
    ])

    # Slide 13: Ensemble Strategy
    add_slide("Ensemble Strategy - Combining Models", [
        "Approach: Stacking Ensemble",
        "",
        "Level 1 (Base Models):",
        "  XGBoost --> P(winner=1) for each ball",
        "  Random Forest --> P(winner=1) for each ball",
        "  LightGBM --> P(winner=1) for each ball",
        "  Neural Net --> P(winner=1) for each ball",
        "",
        "Level 2 (Meta-Learner):",
        "  Logistic Regression on Level-1 predictions",
        "  Input: [xgb_prob, rf_prob, lgbm_prob, nn_prob]",
        "  Output: final P(winner=1)",
        "",
        "Final Selection:",
        "  Rank all balls by ensemble probability",
        "  Select top-N balls --> apply best_practice_state filter",
        "  Output: 5 prediction lines matching LotteryOptimizerEngine pred5 format",
    ])

    # Slide 14: Integration with Existing System
    add_slide("Integration with LotteryOptimizerEngine.js", [
        "Current System (rule-based):",
        "  20 strategies --> checkCombination() --> score --> pred5",
        "",
        "ML-Enhanced System (hybrid):",
        "  1. Extract features using existing strategy functions",
        "  2. Feed features to trained ML models",
        "  3. Get probability rankings from ensemble",
        "  4. Filter through best_practice_state (existing strategies)",
        "  5. Output hybrid pred5: ML-ranked + strategy-validated",
        "",
        "API Endpoint (proposed):",
        "  POST /ml/predict",
        "  Body: {state_code, game_type, draw_time, model_id}",
        "  Response: {predictions: [{ball, probability, rank, star_rating}]}",
        "",
        "Feedback Loop: PredictionTracker --> retrain models weekly",
    ])

    # Slide 15: Quick Start
    add_slide("Quick Start - Get Running in 30 Minutes", [
        "1. Extract training data:",
        "   curl API/extract?state_code=CA&lottery_ids=powerball&from_date=2020-01-01",
        "",
        "2. Generate features:",
        "   python ml_models/generate_ml_artifacts.py",
        "",
        "3. Train your first model:",
        "   See ml_models/example_xgboost_training.py",
        "",
        "4. Evaluate:",
        "   python ml_models/evaluate_model.py --model xgb_powerball_v1",
        "",
        "5. Predict next draw:",
        "   python ml_models/predict_next_draw.py --game powerball --state CA",
        "",
        "Full guide: ML_LOTTERY_MODEL_GUIDE.md",
    ])

    prs.save(filepath)
    print(f"  PPTX: {filepath}")


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    out_dir = os.path.dirname(os.path.abspath(__file__))
    print("Generating ML artifacts for Lottery Ball State Prediction...")
    print()

    # 1. CSV
    csv_path = os.path.join(out_dir, "ml_training_data.csv")
    generate_example_csv(csv_path, num_draws=200)

    # 2. JSON Schema
    json_path = os.path.join(out_dir, "ml_training_schema.json")
    generate_json_schema(json_path)

    # 3. SQL Schema
    sql_path = os.path.join(out_dir, "ml_database_schema.sql")
    generate_sql_schema(sql_path)

    # 4. PowerPoint
    pptx_path = os.path.join(out_dir, "ml_lottery_pipeline_deck.pptx")
    generate_pptx(pptx_path)

    print()
    print("All ML artifacts generated successfully!")
    print(f"  Output directory: {out_dir}")
